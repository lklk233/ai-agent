# -*- coding: utf-8 -*-
"""文档解析：PDF（文字层/OCR）→ DOCX → TXT/MD → 其他格式兜底。

- .pdf: PyMuPDF(fitz) 提取文字层；页字符数 < min_chars 则渲染图片 + Tesseract OCR
  （chi_sim 中文；OCR 结果按文件哈希缓存到 ocr_cache_dir）
- .docx/.doc: python-docx（.doc 老格式走 textract 兜底）
- .txt/.md: chardet 编码检测
- .pptx: python-pptx；其他: textract 兜底
- 失败文件返回 None 并由调用方记录跳过
"""
import logging
import shutil
import subprocess
import time
from pathlib import Path

logger = logging.getLogger("arch.extractor")


def _tesseract_available(cfg) -> bool:
    cmd = cfg.get("tesseract_cmd", "tesseract")
    return shutil.which(cmd) is not None or Path(cmd).exists()


def extract_pdf_text(path: Path, cfg) -> str:
    """PDF 文字层提取；无文字层则 OCR。返回 (text, method)。"""
    import fitz  # PyMuPDF

    doc = fitz.open(str(path))
    try:
        pages_text = []
        total_chars = 0
        for page in doc:
            t = page.get_text("text") or ""
            pages_text.append(t)
            total_chars += len(t.strip())
        min_chars = int(cfg.get("min_chars_for_text_layer", 10)) * len(doc)
        if total_chars >= min_chars and total_chars > 50:
            return "\n".join(pages_text), "pdf-text"
        # 文字层不足 → OCR
        if not _tesseract_available(cfg):
            logger.warning("%s 无文字层且 Tesseract 不可用，跳过 OCR",
                           path.name)
            return "\n".join(pages_text), "pdf-text-sparse"
        ocr_text = _ocr_pdf(doc, path, cfg)
        return ocr_text, "pdf-ocr"
    finally:
        doc.close()


def _ocr_pdf(doc, path: Path, cfg) -> str:
    """整本 PDF 渲染为图片后 Tesseract 识别；按文件哈希缓存。"""
    import hashlib

    cache_dir = Path(cfg["ocr_cache_dir"])
    cache_dir.mkdir(parents=True, exist_ok=True)
    h = hashlib.sha256(path.read_bytes()).hexdigest()[:16]
    cache_file = cache_dir / f"{h}.txt"
    if cache_file.exists():
        logger.info("%s OCR 命中缓存", path.name)
        return cache_file.read_text(encoding="utf-8", errors="ignore")

    tesseract = cfg.get("tesseract_cmd", "tesseract")
    lang = cfg.get("ocr_lang", "chi_sim+eng")
    dpi = cfg.get("ocr_dpi", 200)
    texts = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(dpi=dpi)
        img_bytes = pix.tobytes("png")
        # 写入临时图片供 tesseract 读取
        tmp_img = cache_dir / f"_tmp_{h}_{i}.png"
        tmp_img.write_bytes(img_bytes)
        try:
            r = subprocess.run(
                [tesseract, str(tmp_img), "stdout", "-l", lang],
                capture_output=True, timeout=120)
            texts.append(r.stdout.decode("utf-8", errors="ignore"))
        except subprocess.TimeoutExpired:
            logger.warning("%s 第 %d 页 OCR 超时", path.name, i)
        finally:
            try:
                tmp_img.unlink()
            except OSError:
                pass
        if i >= int(cfg.get("ocr_max_pages", 100)):
            logger.warning("%s 超过 OCR 页数上限，截断", path.name)
            break
    out = "\n".join(texts)
    cache_file.write_text(out, encoding="utf-8")
    return out


def extract_docx_text(path: Path) -> str:
    import docx
    d = docx.Document(str(path))
    parts = [p.text for p in d.paragraphs if p.text.strip()]
    for tbl in d.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_txt_text(path: Path) -> str:
    raw = path.read_bytes()
    try:
        import chardet
        enc = chardet.detect(raw)["encoding"] or "utf-8"
    except Exception:
        enc = "utf-8"
    for e in (enc, "utf-8", "gb18030", "utf-16"):
        try:
            return raw.decode(e)
        except (UnicodeDecodeError, LookupError):
            continue
    return raw.decode("utf-8", errors="ignore")


def extract_file(path: Path, cfg) -> tuple:
    """统一入口。返回 (text, method) 或 (None, reason)。"""
    path = Path(path)
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return extract_pdf_text(path, cfg)
        if ext in (".docx", ".docm"):
            return extract_docx_text(path), "docx"
        if ext in (".txt", ".md", ".markdown", ".csv"):
            return extract_txt_text(path), "txt"
        if ext in (".pptx", ".ppt"):
            import pptx
            prs = pptx.Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts), "pptx"
        # 其他格式（.doc/.rtf/.msg 等）textract 兜底
        import textract
        return textract.process(str(path)).decode("utf-8", errors="ignore"), "textract"
    except Exception as e:  # noqa: BLE001
        logger.error("解析 %s 失败: %s", path.name, e)
        return None, str(e)


def scan_corpus(corpus_dir: Path):
    """扫描语料目录，返回支持扩展名的文件列表。"""
    exts = {".pdf", ".docx", ".docm", ".doc", ".txt", ".md", ".markdown",
            ".pptx", ".ppt", ".rtf"}
    return sorted(p for p in Path(corpus_dir).rglob("*")
                  if p.is_file() and p.suffix.lower() in exts
                  and not p.name.startswith("~$"))
